#!/usr/bin/env python3
"""
分析PPT文件结构，查看幻灯片内容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

def analyze_ppt(file_path):
    """分析PPT文件"""
    print(f"分析PPT文件: {file_path}")

    # 打开PPT文件
    prs = Presentation(file_path)

    # 获取幻灯片数量
    slide_count = len(prs.slides)
    print(f"幻灯片数量: {slide_count}")

    # 遍历每张幻灯片
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n=== 幻灯片 {i} ===")

        # 获取幻灯片中的形状
        shape_count = len(slide.shapes)
        print(f"形状数量: {shape_count}")

        # 遍历每个形状
        for j, shape in enumerate(slide.shapes, 1):
            print(f"  形状 {j}: {shape.shape_type}, 名称: {shape.name}")

            # 如果是文本框，显示文本内容
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    # 处理Unicode字符
                    try:
                        print(f"    文本: {text[:100]}...")
                    except UnicodeEncodeError:
                        # 替换无法编码的字符
                        safe_text = text.encode('utf-8', errors='replace').decode('utf-8')
                        print(f"    文本: {safe_text[:100]}...")

        # 检查是否有动画
        if slide.has_notes_slide:
            print("  有备注幻灯片")

        # 检查过渡效果（使用底层XML）
        try:
            transition = slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
            if transition is not None:
                print("  [OK] 有过渡效果")
            else:
                print("  [NO] 无过渡效果")
        except Exception as e:
            print(f"  [ERR] 检查过渡效果时出错: {e}")

if __name__ == "__main__":
    file_path = "dyn-plant-echo_20260528_203622.pptx"
    if os.path.exists(file_path):
        analyze_ppt(file_path)
    else:
        print(f"文件不存在: {file_path}")