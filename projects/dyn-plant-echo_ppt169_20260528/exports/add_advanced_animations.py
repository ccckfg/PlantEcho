#!/usr/bin/env python3
"""
为PPT添加高级动画效果
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

def add_appear_animation(slide, shape, delay=0):
    """为形状添加出现动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:set>
                                            <p:cBhvr>
                                                <p:cTn id="5" dur="1" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                                <p:attrNameLst>
                                                    <p:attrName>style.visibility</p:attrName>
                                                </p:attrNameLst>
                                            </p:cBhvr>
                                            <p:to>
                                                <p:strVal val="visible"/>
                                            </p:to>
                                        </p:set>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_fade_animation(slide, shape, delay=0, duration=500):
    """为形状添加淡入动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:anim>
                                            <p:cBhvr additive="base">
                                                <p:cTn id="5" dur="{duration}" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                                <p:attrNameLst>
                                                    <p:attrName>style.opacity</p:attrName>
                                                </p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst>
                                                <p:tav tm="0">
                                                    <p:val>
                                                        <p:strVal val="0"/>
                                                    </p:val>
                                                </p:tav>
                                                <p:tav tm="100000">
                                                    <p:val>
                                                        <p:strVal val="1"/>
                                                    </p:val>
                                                </p:tav>
                                            </p:tavLst>
                                        </p:anim>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_scale_animation(slide, shape, delay=0, duration=500, from_scale=0.5, to_scale=1.0):
    """为形状添加缩放动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:animScale>
                                            <p:cBhvr>
                                                <p:cTn id="5" dur="{duration}" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                            </p:cBhvr>
                                            <p:from x="{from_scale * 100}%" y="{from_scale * 100}%"/>
                                            <p:to x="{to_scale * 100}%" y="{to_scale * 100}%"/>
                                        </p:animScale>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_rotate_animation(slide, shape, delay=0, duration=500, from_angle=0, to_angle=360):
    """为形状添加旋转动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:animRot>
                                            <p:cBhvr>
                                                <p:cTn id="5" dur="{duration}" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                            </p:cBhvr>
                                            <p:from ang="{from_angle * 60000}"/>
                                            <p:to ang="{to_angle * 60000}"/>
                                        </p:animRot>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_advanced_animations_to_ppt(input_file, output_file):
    """为PPT添加高级动画"""
    print(f"正在处理PPT文件: {input_file}")

    # 打开PPT文件
    prs = Presentation(input_file)

    # 遍历每张幻灯片
    for i, slide in enumerate(prs.slides, 1):
        print(f"正在处理幻灯片 {i}...")

        # 为每个形状添加动画
        for j, shape in enumerate(slide.shapes, 1):
            # 为文本框添加淡入动画
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and not text.startswith("0") and not text.startswith("/"):  # 跳过页码
                    print(f"  为形状 {j} ({shape.name}) 添加淡入动画")
                    add_fade_animation(slide, shape, delay=j * 200, duration=500)

            # 为椭圆形添加缩放动画
            elif shape.shape_type == 1 and "Ellipse" in shape.name:
                print(f"  为形状 {j} ({shape.name}) 添加缩放动画")
                add_scale_animation(slide, shape, delay=j * 100, duration=300, from_scale=0.5, to_scale=1.0)

            # 为矩形添加出现动画
            elif shape.shape_type == 1 and "Rectangle" in shape.name:
                print(f"  为形状 {j} ({shape.name}) 添加出现动画")
                add_appear_animation(slide, shape, delay=j * 100)

    # 保存PPT文件
    print(f"正在保存PPT文件: {output_file}")
    prs.save(output_file)
    print("高级动画添加完成！")

if __name__ == "__main__":
    input_file = "dyn-plant-echo_20260528_203622.pptx"
    output_file = "dyn-plant-echo_20260528_203622_with_advanced_animations.pptx"

    if os.path.exists(input_file):
        add_advanced_animations_to_ppt(input_file, output_file)
    else:
        print(f"文件不存在: {input_file}")