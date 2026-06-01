#!/usr/bin/env python3
"""
为PPT添加动画效果
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

def add_appear_animation(slide, shape, delay=0):
    """为形状添加出现动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml(f'<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:set>
                                            <p:cBhvr>
                                                <p:cTn id="5" dur="1" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                                <p:attrNameLst>
                                                    <p:attrName>style.visibility</p:attrName>
                                                </p:attrNameLst>
                                            </p:cBhvr>
                                            <p:to>
                                                <p:strVal val="visible"/>
                                            </p:to>
                                        </p:set>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_fly_in_animation(slide, shape, delay=0, direction="fromLeft"):
    """为形状添加飞入动画"""
    # 获取或创建timing元素
    timing = slide._element.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        slide._element.append(timing)

    # 获取或创建tnLst元素
    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
        timing.append(tnLst)

    # 获取或创建par元素
    par = tnLst.find(qn('p:par'))
    if par is None:
        par = parse_xml('<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst/></p:cTn></p:par>')
        tnLst.append(par)

    # 获取childTnLst
    childTnLst = par.find(qn('p:cTn')).find(qn('p:childTnLst'))

    # 创建动画序列
    seq_xml = f'''
    <p:seq xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" concurrent="1" nextAc="seek">
        <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
            <p:childTnLst>
                <p:par>
                    <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                            <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                            <p:par>
                                <p:cTn id="4" fill="hold">
                                    <p:stCondLst>
                                        <p:cond delay="{delay}"/>
                                    </p:stCondLst>
                                    <p:childTnLst>
                                        <p:anim>
                                            <p:cBhvr additive="base">
                                                <p:cTn id="5" dur="500" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                </p:cTn>
                                                <p:tgtEl>
                                                    <p:spTgt spid="{shape.shape_id}"/>
                                                </p:tgtEl>
                                                <p:attrNameLst>
                                                    <p:attrName>ppt_x</p:attrName>
                                                    <p:attrName>ppt_y</p:attrName>
                                                </p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst>
                                                <p:tav tm="0">
                                                    <p:val>
                                                        <p:strVal val="#ppt_x"/>
                                                    </p:val>
                                                </p:tav>
                                                <p:tav tm="100000">
                                                    <p:val>
                                                        <p:strVal val="0"/>
                                                    </p:val>
                                                </p:tav>
                                            </p:tavLst>
                                        </p:anim>
                                    </p:childTnLst>
                                </p:cTn>
                            </p:par>
                        </p:childTnLst>
                    </p:cTn>
                </p:par>
            </p:childTnLst>
        </p:cTn>
        <p:prevCondLst>
            <p:cond evt="onPrev" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:prevCondLst>
        <p:nextCondLst>
            <p:cond evt="onNext" delay="0">
                <p:tgtEl>
                    <p:sldTgt/>
                </p:tgtEl>
            </p:cond>
        </p:nextCondLst>
    </p:seq>
    '''

    seq = parse_xml(seq_xml)
    childTnLst.append(seq)

    return True

def add_animations_to_ppt(input_file, output_file):
    """为PPT添加动画"""
    print(f"正在处理PPT文件: {input_file}")

    # 打开PPT文件
    prs = Presentation(input_file)

    # 遍历每张幻灯片
    for i, slide in enumerate(prs.slides, 1):
        print(f"正在处理幻灯片 {i}...")

        # 为每个形状添加动画
        for j, shape in enumerate(slide.shapes, 1):
            # 为文本框添加出现动画
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and not text.startswith("0") and not text.startswith("/"):  # 跳过页码
                    print(f"  为形状 {j} ({shape.name}) 添加出现动画")
                    add_appear_animation(slide, shape, delay=j * 200)

            # 为其他形状添加动画
            elif shape.shape_type == 1:  # AUTO_SHAPE
                print(f"  为形状 {j} ({shape.name}) 添加出现动画")
                add_appear_animation(slide, shape, delay=j * 100)

    # 保存PPT文件
    print(f"正在保存PPT文件: {output_file}")
    prs.save(output_file)
    print("动画添加完成！")

if __name__ == "__main__":
    input_file = "dyn-plant-echo_20260528_203622.pptx"
    output_file = "dyn-plant-echo_20260528_203622_with_animations.pptx"

    if os.path.exists(input_file):
        add_animations_to_ppt(input_file, output_file)
    else:
        print(f"文件不存在: {input_file}")